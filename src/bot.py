import os
import tempfile
import json
from typing import BinaryIO, Dict, Any
from loguru import logger
from telegram import Update, InlineKeyboardButton, InlineKeyboardMarkup
from telegram.ext import (
    Application,
    CommandHandler,
    MessageHandler,
    CallbackQueryHandler,
    ConversationHandler,
    ContextTypes,
    filters
)
from services.rag_service import rag_service
from services.database_service import database_service
import PyPDF2
import docx
import openpyxl
from pptx import Presentation

# Conversation states
AWAITING_INPUT = 1

# Configure logging
logger.add("bot.log", rotation="500 MB")

async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Send a message when the command /start is issued."""
    keyboard = [
        [InlineKeyboardButton("⚙️ Налаштування", callback_data='settings')],
        [InlineKeyboardButton("👥 Клієнти", callback_data='clients')]
    ]
    reply_markup = InlineKeyboardMarkup(keyboard)
    
    await update.message.reply_text(
        "👋 Привіт! Я твій AI-асистент в OVB. Ти можеш:\n"
        "1. Надсилати мені голосові повідомлення з питаннями\n"
        "2. Надсилати документи для навчання\n"
        "3. Керувати базою клієнтів\n"
        "Я використаю свої знання, щоб відповісти на ваші запитання!",
        reply_markup=reply_markup
    )

async def help_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Send a message when the command /help is issued."""
    await update.message.reply_text(
        "Як мене використовувати:\n\n"
        "🎤 Голосові повідомлення: Задавайте питання голосовим повідомленням\n"
        "📄 Документи: Надсилайте текстові файли для навчання\n"
        "❓ Питання: Я використаю вивчене, щоб відповісти на ваші запитання\n\n"
        "Підтримувані формати файлів:\n"
        "- Текстові файли (.txt, .md)\n"
        "- PDF документи (.pdf)\n"
        "- Word документи (.doc, .docx)\n"
        "- Excel файли (.xlsx)\n"
        "- PowerPoint презентації (.pptx)"
    )

async def settings(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle settings button click."""
    query = update.callback_query
    await query.answer()
    
    keyboard = [
        [InlineKeyboardButton("📚 Додати документи для навчання", callback_data='add_docs')],
        [InlineKeyboardButton("❌ Видалити всі документи", callback_data='delete_docs')],
        [InlineKeyboardButton("ℹ️ Статистика бази знань", callback_data='stats')]
    ]
    reply_markup = InlineKeyboardMarkup(keyboard)
    
    await query.edit_message_text(
        text="⚙️ Налаштування\n\n"
        "Оберіть опцію:",
        reply_markup=reply_markup
    )

async def handle_settings_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle settings menu callbacks."""
    query = update.callback_query
    await query.answer()
    
    if query.data == 'add_docs':
        await query.edit_message_text(
            "📚 Надішліть мені документи для навчання.\n\n"
            "Підтримувані формати:\n"
            "- Текстові файли (.txt, .md)\n"
            "- PDF документи (.pdf)\n"
            "- Word документи (.doc, .docx)\n"
            "- Excel файли (.xlsx)\n"
            "- PowerPoint презентації (.pptx)"
        )
    elif query.data == 'delete_docs':
        # Тут можна додати логіку видалення документів
        await query.edit_message_text("🗑 Всі документи видалено з бази знань.")
    elif query.data == 'stats':
        # Тут можна додати логіку отримання статистики
        await query.edit_message_text("📊 Статистика бази знань буде доступна незабаром.")

def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF file."""
    text = ""
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            text += page.extract_text() + "\n"
    return text

def extract_text_from_docx(file_path: str) -> str:
    """Extract text from DOCX file."""
    doc = docx.Document(file_path)
    return "\n".join([paragraph.text for paragraph in doc.paragraphs])

def extract_text_from_xlsx(file_path: str) -> str:
    """Extract text from XLSX file."""
    wb = openpyxl.load_workbook(file_path)
    text = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        for row in ws.rows:
            row_text = " ".join(str(cell.value) for cell in row if cell.value)
            if row_text:
                text.append(row_text)
    return "\n".join(text)

def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PPTX file."""
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

async def handle_voice(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle voice messages."""
    try:
        # Send initial status
        status_message = await update.message.reply_text("🎧 Обробляю ваше голосове повідомлення...")
        
        # Get voice file
        voice = update.message.voice
        voice_file = await context.bot.get_file(voice.file_id)
        
        # Download voice file to temporary location
        with tempfile.NamedTemporaryFile(suffix=".ogg", delete=False) as temp_file:
            await voice_file.download_to_drive(temp_file.name)
            
            # Process audio using RAG service
            result = await rag_service.process_audio_query(temp_file.name)
            
            # Extract client info from the response
            client_info = extract_client_info(result['response'])
            
            # Store client info in context for later use
            if client_info:
                context.user_data['current_client_info'] = client_info
                logger.info(f"Extracted client info: {client_info}")
            
            # Add buttons for saving/editing client info
            keyboard = [
                [
                    InlineKeyboardButton("💾 Зберегти дані клієнта", callback_data='save_client'),
                    InlineKeyboardButton("✏️ Редагувати дані", callback_data='edit_client')
                ]
            ]
            reply_markup = InlineKeyboardMarkup(keyboard)
            
            # Send transcription and response
            await status_message.edit_text(
                f"🎯 Я почув: {result['transcription']}\n\n"
                f"🤖 Моя відповідь: {result['response']}",
                reply_markup=reply_markup
            )
            
        # Cleanup temporary file
        os.unlink(temp_file.name)
        
    except Exception as e:
        logger.error(f"Error processing voice message: {str(e)}")
        await update.message.reply_text(
            "😕 Вибачте, виникла проблема з обробкою голосового повідомлення. "
            "Будь ласка, спробуйте ще раз або переконайтеся, що повідомлення записано чітко."
        )

def extract_client_info(response_text: str) -> Dict[str, Any]:
    """Extract client information from the response text"""
    try:
        client_info = {}
        lines = response_text.split('\n')
        
        # Initialize description collection
        description_lines = []
        collecting_description = False
        
        for line in lines:
            line = line.strip()
            
            # Skip empty lines
            if not line:
                continue
                
            # Handle client name
            if any(x in line.lower() for x in ['клієнт:', 'імʼя:', "ім'я:", 'імя:']):
                client_info['full_name'] = line.split(':', 1)[1].strip()
            
            # Handle age
            elif 'вік:' in line.lower():
                age_text = line.split(':', 1)[1].strip()
                try:
                    client_info['age'] = int(''.join(filter(str.isdigit, age_text)))
                except ValueError:
                    logger.warning(f"Could not parse age from: {age_text}")
            
            # Handle meeting date
            elif 'дата:' in line.lower():
                client_info['meeting_date'] = line.split(':', 1)[1].strip()
            
            # Handle product type
            elif 'продукт:' in line.lower():
                client_info['product_type'] = line.split(':', 1)[1].strip()
            
            # Handle goal
            elif 'ціль:' in line.lower():
                client_info['goal'] = line.split(':', 1)[1].strip()
            
            # Start collecting description
            elif any(x in line.lower() for x in ['опис:', 'опис клієнта:', 'додаткова інформація:']):
                collecting_description = True
                continue
            
            # Stop collecting description if we hit a section marker
            elif line.startswith(('###', '---', '🎯', '💡', '📝')):
                collecting_description = False
            
            # Collect description lines
            elif collecting_description:
                description_lines.append(line)
        
        # Join description lines if we collected any
        if description_lines:
            client_info['description'] = ' '.join(description_lines).strip()
        
        # Validate required fields
        required_fields = ['full_name', 'age', 'product_type', 'goal']
        if all(field in client_info for field in required_fields):
            logger.info(f"Successfully extracted client info: {client_info}")
            return client_info
            
        missing_fields = [field for field in required_fields if field not in client_info]
        logger.warning(f"Missing required fields: {missing_fields}")
        return None
        
    except Exception as e:
        logger.error(f"Error extracting client info: {str(e)}")
        return None

async def handle_clients(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle clients button click"""
    query = update.callback_query
    await query.answer()
    
    try:
        # Get all clients from database
        clients = await database_service.get_all_clients()
        
        if not clients:
            await query.edit_message_text(
                "📝 Список клієнтів порожній. Додайте клієнтів через голосові повідомлення."
            )
            return
        
        # Format client list
        client_text = "📋 Список клієнтів:\n\n"
        for client in clients:
            client_text += (
                f"👤 {client['full_name']}\n"
                f"📅 Вік: {client['age']}\n"
                f"🎯 Ціль: {client['goal']}\n"
                f"💼 Продукт: {client['product_type']}\n"
                f"📝 Опис: {client['description'][:100]}...\n\n"
            )
        
        # Add navigation buttons if needed
        keyboard = [[InlineKeyboardButton("🔄 Оновити", callback_data='clients')]]
        reply_markup = InlineKeyboardMarkup(keyboard)
        
        await query.edit_message_text(
            client_text,
            reply_markup=reply_markup
        )
        
    except Exception as e:
        logger.error(f"Error getting clients: {str(e)}")
        await query.edit_message_text(
            "❌ Помилка при отриманні списку клієнтів. Спробуйте пізніше."
        )

async def handle_client_action(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle client-related button actions"""
    query = update.callback_query
    await query.answer()
    
    if query.data == 'save_client':
        client_info = context.user_data.get('current_client_info')
        if not client_info:
            await query.message.reply_text(
                "❌ Немає даних для збереження. Спочатку надішліть голосове повідомлення."
            )
            return
            
        try:
            # Validate required fields
            required_fields = ['full_name', 'age', 'product_type', 'goal']
            missing_fields = [field for field in required_fields if field not in client_info]
            
            if missing_fields:
                await query.message.reply_text(
                    f"❌ Відсутні обов'язкові поля: {', '.join(missing_fields)}.\n"
                    "Будь ласка, відредагуйте дані клієнта."
                )
                return
            
            # Save client data
            client_id = await database_service.save_client(client_info)
            
            # Send success message with client details
            success_message = (
                f"✅ Дані клієнта успішно збережено!\n\n"
                f"📋 Інформація про клієнта:\n"
                f"🆔 ID: {client_id}\n"
                f"👤 Ім'я: {client_info['full_name']}\n"
                f"📅 Вік: {client_info['age']}\n"
                f"💼 Продукт: {client_info['product_type']}\n"
                f"🎯 Ціль: {client_info['goal']}\n"
            )
            
            if client_info.get('description'):
                success_message += f"📝 Опис: {client_info['description']}\n"
            
            keyboard = [
                [
                    InlineKeyboardButton("✏️ Редагувати", callback_data='edit_saved_client'),
                    InlineKeyboardButton("👥 Всі клієнти", callback_data='clients')
                ]
            ]
            reply_markup = InlineKeyboardMarkup(keyboard)
            
            # Send success message as a new message
            await query.message.reply_text(success_message, reply_markup=reply_markup)
            
            # Also update the original message to show it was processed
            await query.edit_message_text(
                query.message.text + "\n\n✅ Дані успішно збережено!"
            )
            
            # Clear the stored client info after successful save
            context.user_data['current_client_info'] = None
            
        except Exception as e:
            logger.error(f"Error saving client: {str(e)}")
            await query.message.reply_text(
                "❌ Помилка при збереженні даних клієнта. Спробуйте ще раз."
            )
    
    elif query.data == 'edit_client' or query.data == 'edit_saved_client':
        client_info = context.user_data.get('current_client_info')
        if not client_info:
            await query.message.reply_text(
                "❌ Немає даних для редагування. Спочатку надішліть голосове повідомлення."
            )
            return
        
        # Show current values and edit options
        current_values = (
            f"📋 Поточні дані клієнта:\n\n"
            f"👤 Ім'я: {client_info.get('full_name', 'Не вказано')}\n"
            f"📅 Вік: {client_info.get('age', 'Не вказано')}\n"
            f"🎯 Ціль: {client_info.get('goal', 'Не вказано')}\n"
            f"💼 Продукт: {client_info.get('product_type', 'Не вказано')}\n"
            f"📝 Опис: {client_info.get('description', 'Не вказано')}\n\n"
            "Оберіть поле для редагування:"
        )
        
        edit_keyboard = [
            [
                InlineKeyboardButton("👤 Ім'я", callback_data='edit_name'),
                InlineKeyboardButton("📅 Вік", callback_data='edit_age')
            ],
            [
                InlineKeyboardButton("🎯 Ціль", callback_data='edit_goal'),
                InlineKeyboardButton("💼 Продукт", callback_data='edit_product')
            ],
            [InlineKeyboardButton("📝 Опис", callback_data='edit_description')],
            [
                InlineKeyboardButton("💾 Зберегти", callback_data='save_changes'),
                InlineKeyboardButton("❌ Скасувати", callback_data='cancel_edit')
            ]
        ]
        reply_markup = InlineKeyboardMarkup(edit_keyboard)
        
        await query.message.reply_text(current_values, reply_markup=reply_markup)

async def handle_edit_field(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle editing specific client fields"""
    query = update.callback_query
    await query.answer()
    
    field_map = {
        'edit_name': ('full_name', "👤 Введіть нове ім'я клієнта:"),
        'edit_age': ('age', "📅 Введіть новий вік клієнта:"),
        'edit_goal': ('goal', "🎯 Введіть нову ціль клієнта:"),
        'edit_product': ('product_type', "💼 Введіть новий тип продукту:"),
        'edit_description': ('description', "📝 Введіть новий опис клієнта:")
    }
    
    if query.data.startswith('edit_'):
        field = query.data.replace('edit_', '')
        if field in [k.replace('edit_', '') for k in field_map.keys()]:
            field_name, prompt = field_map[f'edit_{field}']
            context.user_data['editing_field'] = field_name
            
            # Show current value if exists
            client_info = context.user_data.get('current_client_info', {})
            current_value = client_info.get(field_name, 'Не вказано')
            
            keyboard = [[
                InlineKeyboardButton("❌ Скасувати", callback_data='cancel_edit')
            ]]
            reply_markup = InlineKeyboardMarkup(keyboard)
            
            await query.edit_message_text(
                f"{prompt}\n\n"
                f"Поточне значення: {current_value}\n\n"
                "Надішліть нове значення як текстове повідомлення",
                reply_markup=reply_markup
            )
            return AWAITING_INPUT
    
    elif query.data == 'cancel_edit':
        keyboard = [
            [InlineKeyboardButton("↩️ Повернутися до редагування", callback_data='edit_client')]
        ]
        reply_markup = InlineKeyboardMarkup(keyboard)
        
        await query.edit_message_text(
            "❌ Редагування скасовано",
            reply_markup=reply_markup
        )
        return ConversationHandler.END
    
    return ConversationHandler.END

async def handle_text_input(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle text input for editing"""
    field = context.user_data.get('editing_field')
    if not field:
        return ConversationHandler.END
    
    text = update.message.text
    client_info = context.user_data.get('current_client_info', {})
    
    if field == 'age':
        try:
            client_info[field] = int(text)
        except ValueError:
            await update.message.reply_text(
                "❌ Будь ласка, введіть коректний вік (число)",
                reply_markup=InlineKeyboardMarkup([[
                    InlineKeyboardButton("↩️ Спробувати ще раз", callback_data=f'edit_age')
                ]])
            )
            return AWAITING_INPUT
    else:
        client_info[field] = text
    
    context.user_data['current_client_info'] = client_info
    
    keyboard = [
        [
            InlineKeyboardButton("✏️ Редагувати ще", callback_data='edit_client'),
            InlineKeyboardButton("💾 Зберегти", callback_data='save_changes')
        ]
    ]
    reply_markup = InlineKeyboardMarkup(keyboard)
    
    await update.message.reply_text(
        f"✅ Значення оновлено успішно!\n"
        f"Поле: {field}\n"
        f"Нове значення: {text}",
        reply_markup=reply_markup
    )
    
    return ConversationHandler.END

async def handle_document(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle document messages."""
    try:
        # Check if document format is supported
        document = update.message.document
        file_name = document.file_name.lower()
        supported_formats = ('.txt', '.md', '.pdf', '.doc', '.docx', '.xlsx', '.pptx')
        
        if not any(file_name.endswith(fmt) for fmt in supported_formats):
            await update.message.reply_text(
                "❌ Будь ласка, надсилайте файли лише в підтримуваних форматах:\n"
                "- Текстові файли (.txt, .md)\n"
                "- PDF документи (.pdf)\n"
                "- Word документи (.doc, .docx)\n"
                "- Excel файли (.xlsx)\n"
                "- PowerPoint презентації (.pptx)"
            )
            return
            
        # Send initial status
        status_message = await update.message.reply_text("📚 Обробляю ваш документ...")
        
        # Get document file
        doc_file = await context.bot.get_file(document.file_id)
        
        # Download and process document
        with tempfile.NamedTemporaryFile(suffix=os.path.splitext(file_name)[1], delete=False) as temp_file:
            await doc_file.download_to_drive(temp_file.name)
            
            # Extract text based on file format
            if file_name.endswith('.pdf'):
                text = extract_text_from_pdf(temp_file.name)
            elif file_name.endswith(('.doc', '.docx')):
                text = extract_text_from_docx(temp_file.name)
            elif file_name.endswith('.xlsx'):
                text = extract_text_from_xlsx(temp_file.name)
            elif file_name.endswith('.pptx'):
                text = extract_text_from_pptx(temp_file.name)
            else:  # txt or md
                with open(temp_file.name, 'r', encoding='utf-8') as f:
                    text = f.read()
            
            # Process document using RAG service
            await rag_service.process_document(text, document.file_id)
            
            await status_message.edit_text(
                "✅ Документ успішно оброблено! Я вивчив його вміст."
            )
            
        # Cleanup temporary file
        os.unlink(temp_file.name)
        
    except Exception as e:
        logger.error(f"Error processing document: {str(e)}")
        await update.message.reply_text(
            "😕 Вибачте, виникла проблема з обробкою документа. "
            "Будь ласка, перевірте формат файлу та спробуйте ще раз."
        )

def main():
    """Start the bot."""
    # Create application
    application = Application.builder().token(os.getenv("TELEGRAM_BOT_TOKEN")).build()

    # Add handlers
    application.add_handler(CommandHandler("start", start))
    application.add_handler(CommandHandler("help", help_command))
    application.add_handler(CallbackQueryHandler(settings, pattern='^settings$'))
    application.add_handler(CallbackQueryHandler(handle_settings_callback, pattern='^(add_docs|delete_docs|stats)$'))
    application.add_handler(CallbackQueryHandler(handle_clients, pattern='^clients$'))
    application.add_handler(CallbackQueryHandler(handle_client_action, pattern='^(save_client|edit_client|save_changes)$'))
    
    # Add handler for text messages during editing
    application.add_handler(MessageHandler(
        filters.TEXT & ~filters.COMMAND,
        handle_text_input,
        block=False
    ))
    
    # Add handler for edit field callbacks
    application.add_handler(CallbackQueryHandler(
        handle_edit_field,
        pattern='^(edit_|cancel_edit)'
    ))
    
    application.add_handler(MessageHandler(filters.VOICE, handle_voice))
    application.add_handler(MessageHandler(filters.Document.ALL, handle_document))

    # Start the bot
    application.run_polling(allowed_updates=Update.ALL_TYPES)

if __name__ == '__main__':
    main() 